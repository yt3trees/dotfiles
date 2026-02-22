"""
slide_generator_pptx.py - PowerPoint生成スクリプト

slides.md を読み込み、15レイアウトに対応した PPTX を生成する。
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from slides_markdown import parse_markdown_file, SlideData, char_width, truncate_text

# スライドサイズ 16:9 (EMU)
SLIDE_WIDTH = 12192000
SLIDE_HEIGHT = 6858000


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
class Config:
    """config.json を読み込んで設定値を提供する"""

    def __init__(self, config_path: str | None = None):
        self.data = {}
        if config_path and os.path.exists(config_path):
            with open(config_path, 'r', encoding='utf-8') as f:
                self.data = json.load(f)

        palette = self.data.get('palette', {})
        self.primary = palette.get('primary', '#1E3A5F')
        self.secondary = palette.get('secondary', '#4A6FA5')
        self.accent = palette.get('accent', '#3AA899')
        self.gray = palette.get('gray', '#999999')

        text = palette.get('text', {})
        self.text_primary = text.get('primary', '#333333')
        self.text_secondary = text.get('secondary', '#666666')
        self.text_light = text.get('light', '#FFFFFF')

        bg = palette.get('background', {})
        self.bg_primary = bg.get('primary', '#FFFFFF')
        self.bg_secondary = bg.get('secondary', '#F5F5F5')
        self.bg_dark = bg.get('dark', '#1E3A5F')

        chart_cfg = palette.get('chart', {})
        self.chart_line_color_3 = chart_cfg.get('line_color_3', '#EDB120')

        font_cfg = self.data.get('font', {})
        self.font_family = font_cfg.get('family', 'Meiryo UI')

        template_cfg = self.data.get('template', {})
        self.template_path = template_cfg.get('pptx_path', '')

    def rgb(self, hex_color: str) -> RGBColor:
        """HEX色文字列をRGBColorに変換する"""
        h = hex_color.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Chart Renderer
# ---------------------------------------------------------------------------
class ChartRenderer:
    """グラフ描画クラス"""

    def __init__(self, config: Config):
        self.cfg = config

    def render(self, slide, slide_data: SlideData, left, top, width, height):
        """グラフ種別に応じたグラフを描画する"""
        table = slide_data.table
        if not table or not table.rows:
            return

        chart_type = slide_data.chart_type
        if chart_type == "pie":
            self._render_pie(slide, table, left, top, width, height)
        elif chart_type == "line":
            self._render_line(slide, table, left, top, width, height)
        elif chart_type == "bar":
            self._render_bar(slide, table, left, top, width, height)
        else:
            self._render_horizontal_bar(slide, table, left, top, width, height)

    def _get_labels_and_series(self, table):
        """テーブルからラベルと系列データを抽出する"""
        labels = []
        highlight_flags = []
        for row in table.rows:
            raw_label = row[0]
            is_bold = '**' in raw_label
            clean_label = raw_label.replace('**', '').strip()
            labels.append(clean_label)
            highlight_flags.append(is_bold)

        series_names = table.headers[1:]
        series_data = {name: [] for name in series_names}
        for row in table.rows:
            for i, name in enumerate(series_names):
                val_str = row[i + 1].replace('**', '').replace('%', '').replace(',', '').strip() if i + 1 < len(row) else '0'
                try:
                    series_data[name].append(float(val_str))
                except ValueError:
                    series_data[name].append(0)

        return labels, highlight_flags, series_names, series_data

    def _render_horizontal_bar(self, slide, table, left, top, width, height):
        """横棒グラフ"""
        labels, highlight_flags, series_names, series_data = self._get_labels_and_series(table)

        chart_data = CategoryChartData()
        chart_data.categories = labels
        for name in series_names:
            chart_data.add_series(name, series_data[name])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = False

        # スタイル設定
        plot = chart.plots[0]
        plot.gap_width = 80

        # カテゴリ軸(横棒の場合は値軸が下、カテゴリ軸が左)
        cat_axis = chart.category_axis
        cat_axis.has_major_gridlines = False
        cat_axis.major_tick_mark = 2  # XL_TICK_MARK.NONE
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = self.cfg.font_family
        cat_axis.tick_labels.font.color.rgb = self.cfg.rgb(self.cfg.text_primary)
        cat_axis.format.line.fill.background()

        val_axis = chart.value_axis
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.dash_style = 4  # MSO_LINE_DASH_STYLE.DASH
        val_axis.major_gridlines.format.line.fill.solid()
        val_axis.major_gridlines.format.line.fill.fore_color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        val_axis.major_tick_mark = 2  # NONE
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = self.cfg.font_family
        val_axis.format.line.fill.background()

        # バーの色(強調あり/なし)
        series = plot.series[0]
        for idx, is_highlight in enumerate(highlight_flags):
            point = series.points[idx]
            point.format.fill.solid()
            if is_highlight:
                point.format.fill.fore_color.rgb = self.cfg.rgb(self.cfg.accent)
            else:
                point.format.fill.fore_color.rgb = self.cfg.rgb(self.cfg.gray)
            self._set_rounded_rect(point)

    def _render_bar(self, slide, table, left, top, width, height):
        """縦棒グラフ"""
        labels, highlight_flags, series_names, series_data = self._get_labels_and_series(table)

        chart_data = CategoryChartData()
        chart_data.categories = labels
        for name in series_names:
            chart_data.add_series(name, series_data[name])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = False

        plot = chart.plots[0]
        plot.gap_width = 80

        cat_axis = chart.category_axis
        cat_axis.has_major_gridlines = False
        cat_axis.major_tick_mark = 2
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = self.cfg.font_family
        cat_axis.tick_labels.font.color.rgb = self.cfg.rgb(self.cfg.text_primary)
        cat_axis.format.line.fill.background()

        val_axis = chart.value_axis
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.dash_style = 4
        val_axis.major_gridlines.format.line.fill.solid()
        val_axis.major_gridlines.format.line.fill.fore_color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        val_axis.major_tick_mark = 2
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = self.cfg.font_family
        val_axis.format.line.fill.background()

        series = plot.series[0]
        for idx, is_highlight in enumerate(highlight_flags):
            point = series.points[idx]
            point.format.fill.solid()
            if is_highlight:
                point.format.fill.fore_color.rgb = self.cfg.rgb(self.cfg.accent)
            else:
                point.format.fill.fore_color.rgb = self.cfg.rgb(self.cfg.gray)
            self._set_rounded_rect(point)

    def _render_pie(self, slide, table, left, top, width, height):
        """100%積み上げ棒グラフ(円グラフ代替)"""
        labels, highlight_flags, series_names, series_data = self._get_labels_and_series(table)

        chart_data = CategoryChartData()
        chart_data.categories = ['']  # 単一カテゴリ
        for i, label in enumerate(labels):
            vals = [series_data[series_names[0]][i]] if series_names else [0]
            chart_data.add_series(label, vals)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED_100, left, top, width, height, chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = self.cfg.font_family

        plot = chart.plots[0]
        plot.gap_width = 50

        # 軸を非表示
        chart.category_axis.visible = False
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False

        # 色設定
        colors = [self.cfg.accent, self.cfg.secondary, self.cfg.gray, self.cfg.primary]
        for i, series in enumerate(plot.series):
            series.format.fill.solid()
            color = colors[i % len(colors)]
            series.format.fill.fore_color.rgb = self.cfg.rgb(color)

    def _render_line(self, slide, table, left, top, width, height):
        """折れ線グラフ"""
        labels, _, series_names, series_data = self._get_labels_and_series(table)

        chart_data = CategoryChartData()
        chart_data.categories = labels
        for name in series_names:
            chart_data.add_series(name, series_data[name])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
        )
        chart = chart_frame.chart

        if len(series_names) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = self.cfg.font_family
        else:
            chart.has_legend = False

        cat_axis = chart.category_axis
        cat_axis.has_major_gridlines = False
        cat_axis.major_tick_mark = 2
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = self.cfg.font_family
        cat_axis.format.line.fill.background()

        val_axis = chart.value_axis
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.dash_style = 4
        val_axis.major_gridlines.format.line.fill.solid()
        val_axis.major_gridlines.format.line.fill.fore_color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        val_axis.major_tick_mark = 2
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = self.cfg.font_family
        val_axis.format.line.fill.background()

        # 系列の色とマーカー設定
        line_colors = [self.cfg.accent, self.cfg.secondary, self.cfg.chart_line_color_3, self.cfg.primary]
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            color = line_colors[i % len(line_colors)]
            series.format.line.fill.solid()
            series.format.line.fill.fore_color.rgb = self.cfg.rgb(color)
            series.format.line.width = Pt(2.5)
            series.smooth = False
            self._set_line_marker(series, color)

    def _set_rounded_rect(self, point):
        """バーを角丸四角形にする"""
        try:
            sp_pr = point.format._element
            for old in sp_pr.findall(qn('a:prstGeom')):
                sp_pr.remove(old)
            geom = etree.SubElement(sp_pr, qn('a:prstGeom'))
            geom.set('prst', 'roundRect')
            etree.SubElement(geom, qn('a:avLst'))
        except Exception:
            pass

    def _set_line_marker(self, series, color_hex: str):
        """折れ線マーカーを円形・白枠付きに設定"""
        try:
            marker = series.marker
            marker.style = 8  # XL_MARKER_STYLE.CIRCLE
            marker.size = 8
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = self.cfg.rgb(color_hex)
            marker.format.line.fill.solid()
            marker.format.line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            marker.format.line.width = Pt(1.5)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Presentation Builder
# ---------------------------------------------------------------------------
class PresentationBuilder:
    """15レイアウトの描画を行うビルダー"""

    def __init__(self, config: Config, template_path: str | None = None):
        self.cfg = config
        self.chart_renderer = ChartRenderer(config)

        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            self._remove_all_slides()
        else:
            self.prs = Presentation()

        # スライドサイズ設定 16:9
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _remove_all_slides(self):
        """テンプレートの既存スライドを全削除"""
        sldIdLst = self.prs.presentation.sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            self.prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

    def _add_blank_slide(self):
        """空白スライドを追加する"""
        layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    def _remove_placeholders(self, slide):
        """スライドからプレースホルダを除去"""
        sp_tree = slide.shapes._spTree
        for sp in list(sp_tree):
            if sp.tag.endswith('}sp'):
                ph = sp.find('.//' + qn('p:ph'))
                if ph is not None:
                    sp_tree.remove(sp)

    def _set_background(self, slide, hex_color: str):
        """スライド背景色を設定する"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.cfg.rgb(hex_color)

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=12, font_color=None, bold=False,
                     alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     line_spacing=1.5):
        """テキストボックスを追加する"""
        txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = self.cfg.font_family
        p.font.bold = bold
        if font_color:
            p.font.color.rgb = self.cfg.rgb(font_color)
        p.alignment = alignment
        p.space_after = Pt(0)
        p.space_before = Pt(0)

        # 行間設定
        if line_spacing:
            p_elem = p._p
            pPr = p_elem.find(qn('a:pPr'))
            if pPr is None:
                pPr = etree.SubElement(p_elem, qn('a:pPr'))
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spc_pct = etree.SubElement(lnSpc, qn('a:spcPct'))
            spc_pct.set('val', str(int(line_spacing * 100000)))

        # アンカー設定
        txBox.text_frame.paragraphs[0].alignment = alignment
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if anchor == MSO_ANCHOR.MIDDLE:
            bodyPr.set('anchor', 'ctr')
        elif anchor == MSO_ANCHOR.BOTTOM:
            bodyPr.set('anchor', 'b')

        return txBox

    def _add_multiline_textbox(self, slide, left, top, width, height, lines_data,
                               default_font_size=12, default_color=None,
                               alignment=PP_ALIGN.LEFT, line_spacing=1.5):
        """複数行テキストボックスを追加する。lines_dataは(text, font_size, color, bold)のリスト"""
        txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, (text, font_size, color, bold) in enumerate(lines_data):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size or default_font_size)
            p.font.name = self.cfg.font_family
            p.font.bold = bold
            if color:
                p.font.color.rgb = self.cfg.rgb(color)
            elif default_color:
                p.font.color.rgb = self.cfg.rgb(default_color)
            p.alignment = alignment
            p.space_after = Pt(2)
            p.space_before = Pt(2)

            if line_spacing:
                pPr = p._p.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.SubElement(p._p, qn('a:pPr'))
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spc_pct = etree.SubElement(lnSpc, qn('a:spcPct'))
                spc_pct.set('val', str(int(line_spacing * 100000)))

        return txBox

    def _add_rounded_rect(self, slide, left, top, width, height, fill_color, text="",
                          font_size=12, font_color=None, bold=False):
        """角丸四角形シェイプを追加する"""
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE.ROUNDED_RECTANGLE
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.cfg.rgb(fill_color)
        shape.line.fill.background()

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.name = self.cfg.font_family
            p.font.bold = bold
            if font_color:
                p.font.color.rgb = self.cfg.rgb(font_color)
            p.alignment = PP_ALIGN.CENTER

            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            bodyPr.set('anchor', 'ctr')

        return shape

    def _add_circle(self, slide, left, top, size, fill_color, text="",
                    font_size=12, font_color=None, bold=True):
        """円形シェイプを追加する"""
        shape = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(left), Emu(top), Emu(size), Emu(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.cfg.rgb(fill_color)
        shape.line.fill.background()

        if text:
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.name = self.cfg.font_family
            p.font.bold = bold
            if font_color:
                p.font.color.rgb = self.cfg.rgb(font_color)
            p.alignment = PP_ALIGN.CENTER

            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            bodyPr.set('anchor', 'ctr')

        return shape

    def _add_line(self, slide, x1, y1, x2, y2, color, width_pt=1):
        """直線を追加する"""
        from pptx.shapes.connector import Connector
        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            Emu(x1), Emu(y1), Emu(x2), Emu(y2)
        )
        connector.line.fill.solid()
        connector.line.fill.fore_color.rgb = self.cfg.rgb(color)
        connector.line.width = Pt(width_pt)
        return connector

    # --- レイアウト描画メソッド ---

    def _build_title(self, slide_data: SlideData):
        """表紙スライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        self._set_background(slide, self.cfg.bg_dark)

        # タイトル(中央)
        self._add_textbox(
            slide,
            left=Inches(1).emu, top=Inches(2.2).emu,
            width=Inches(10).emu, height=Inches(1.5).emu,
            text=slide_data.title,
            font_size=44, font_color=self.cfg.text_light, bold=True,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
        )

        # サブタイトル
        if slide_data.key_message:
            self._add_textbox(
                slide,
                left=Inches(1).emu, top=Inches(3.8).emu,
                width=Inches(10).emu, height=Inches(0.8).emu,
                text=slide_data.key_message,
                font_size=16, font_color=self.cfg.text_light,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
            )

        # 下部アクセントライン
        self._add_rounded_rect(
            slide,
            left=Inches(4.5).emu, top=Inches(3.5).emu,
            width=Inches(3).emu, height=Inches(0.05).emu,
            fill_color=self.cfg.accent
        )

        return slide

    def _build_section(self, slide_data: SlideData):
        """セクション区切り"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        self._set_background(slide, self.cfg.bg_secondary)

        self._add_textbox(
            slide,
            left=Inches(1).emu, top=Inches(2.5).emu,
            width=Inches(10).emu, height=Inches(1.2).emu,
            text=slide_data.title,
            font_size=36, font_color=self.cfg.primary, bold=True,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
        )

        # 左アクセントライン
        self._add_rounded_rect(
            slide,
            left=Inches(4.5).emu, top=Inches(3.8).emu,
            width=Inches(3).emu, height=Inches(0.05).emu,
            fill_color=self.cfg.accent
        )

        return slide

    def _build_header(self, slide, slide_data: SlideData):
        """共通ヘッダー(タイトル + キーメッセージ)を描画する"""
        y = Inches(0.4).emu

        # タイトル
        if slide_data.title:
            self._add_textbox(
                slide,
                left=Inches(0.6).emu, top=y,
                width=Inches(10.8).emu, height=Inches(0.5).emu,
                text=slide_data.title,
                font_size=24, font_color=self.cfg.text_primary, bold=True,
                alignment=PP_ALIGN.LEFT, line_spacing=1.0
            )
            y += Inches(0.55).emu

        # キーメッセージ
        if slide_data.key_message:
            self._add_textbox(
                slide,
                left=Inches(0.6).emu, top=y,
                width=Inches(10.8).emu, height=Inches(0.4).emu,
                text=slide_data.key_message,
                font_size=16, font_color=self.cfg.text_secondary,
                alignment=PP_ALIGN.LEFT, line_spacing=1.0
            )
            y += Inches(0.5).emu

        # タイトル下のアクセントライン
        self._add_rounded_rect(
            slide,
            left=Inches(0.6).emu, top=y,
            width=Inches(10.8).emu, height=Inches(0.03).emu,
            fill_color=self.cfg.accent
        )

        return y + Inches(0.3).emu

    def _build_toc(self, slide_data: SlideData):
        """目次スライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        items = slide_data.bullets
        if not items:
            return slide

        item_height = Inches(0.5).emu
        circle_size = Inches(0.35).emu
        start_y = content_top + Inches(0.2).emu

        for i, item in enumerate(items):
            y = start_y + i * (item_height + Inches(0.15).emu)

            # 番号付き丸
            self._add_circle(
                slide,
                left=Inches(1.5).emu, top=y,
                size=circle_size,
                fill_color=self.cfg.accent,
                text=str(i + 1),
                font_size=12, font_color=self.cfg.text_light
            )

            # テキスト
            self._add_textbox(
                slide,
                left=Inches(2.1).emu, top=y,
                width=Inches(8).emu, height=item_height,
                text=item,
                font_size=14, font_color=self.cfg.text_primary,
                anchor=MSO_ANCHOR.MIDDLE
            )

        return slide

    def _build_bullet_points(self, slide_data: SlideData):
        """箇条書きスライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        items = slide_data.bullets
        if not items:
            return slide

        start_y = content_top + Inches(0.2).emu
        item_height = Inches(0.45).emu

        for i, item in enumerate(items):
            y = start_y + i * item_height

            # 丸ドット
            dot_size = Inches(0.12).emu
            self._add_circle(
                slide,
                left=Inches(1.0).emu, top=y + Inches(0.15).emu,
                size=dot_size,
                fill_color=self.cfg.accent,
                text="", font_size=1
            )

            # テキスト
            self._add_textbox(
                slide,
                left=Inches(1.3).emu, top=y,
                width=Inches(9.5).emu, height=item_height,
                text=item,
                font_size=12, font_color=self.cfg.text_primary,
                anchor=MSO_ANCHOR.MIDDLE
            )

        return slide

    def _build_numbered_list(self, slide_data: SlideData):
        """番号付きリスト"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        items = slide_data.numbered_items
        if not items:
            return slide

        start_y = content_top + Inches(0.2).emu
        item_height = Inches(0.8).emu
        circle_size = Inches(0.4).emu

        for i, item in enumerate(items):
            y = start_y + i * (item_height + Inches(0.1).emu)

            # 番号バッジ
            self._add_circle(
                slide,
                left=Inches(0.8).emu, top=y + Inches(0.05).emu,
                size=circle_size,
                fill_color=self.cfg.accent,
                text=str(i + 1),
                font_size=14, font_color=self.cfg.text_light
            )

            # タイトル + 説明
            lines_data = [
                (item['title'], 14, self.cfg.text_primary, True),
            ]
            if item.get('description'):
                lines_data.append((item['description'], 11, self.cfg.text_secondary, False))

            self._add_multiline_textbox(
                slide,
                left=Inches(1.5).emu, top=y,
                width=Inches(9.3).emu, height=item_height,
                lines_data=lines_data,
                line_spacing=1.3
            )

        return slide

    def _build_columns(self, slide_data: SlideData, num_cols: int):
        """カラムレイアウト"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        columns = slide_data.columns[:num_cols]
        if not columns:
            return slide

        margin = Inches(0.6).emu
        gap = Inches(0.3).emu
        total_width = SLIDE_WIDTH - 2 * margin
        col_width = (total_width - (num_cols - 1) * gap) // num_cols
        col_height = SLIDE_HEIGHT - content_top - Inches(0.8).emu
        start_y = content_top + Inches(0.3).emu

        for i, col in enumerate(columns):
            x = margin + i * (col_width + gap)

            # カラムの背景カード
            self._add_rounded_rect(
                slide,
                left=x, top=start_y,
                width=col_width, height=col_height,
                fill_color=self.cfg.bg_secondary
            )

            # ヘッダー背景バー
            self._add_rounded_rect(
                slide,
                left=x, top=start_y,
                width=col_width, height=Inches(0.5).emu,
                fill_color=self.cfg.primary,
                text=col.get('heading', ''),
                font_size=12, font_color=self.cfg.text_light, bold=True
            )

            # 本文
            body_text = col.get('body', '')
            self._add_textbox(
                slide,
                left=x + Inches(0.15).emu,
                top=start_y + Inches(0.65).emu,
                width=col_width - Inches(0.3).emu,
                height=col_height - Inches(0.8).emu,
                text=body_text,
                font_size=11, font_color=self.cfg.text_primary,
                line_spacing=1.4
            )

        return slide

    def _build_metrics(self, slide_data: SlideData):
        """数値・KPIスライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        metrics = slide_data.metrics
        if not metrics:
            return slide

        num = len(metrics)
        margin = Inches(1.0).emu
        gap = Inches(0.4).emu
        total_width = SLIDE_WIDTH - 2 * margin
        card_width = (total_width - (num - 1) * gap) // num
        card_height = Inches(2.2).emu
        start_y = content_top + Inches(0.5).emu

        for i, m in enumerate(metrics):
            x = margin + i * (card_width + gap)

            # カード背景
            self._add_rounded_rect(
                slide, left=x, top=start_y,
                width=card_width, height=card_height,
                fill_color=self.cfg.bg_secondary
            )

            # 数値(大きく)
            self._add_textbox(
                slide,
                left=x, top=start_y + Inches(0.3).emu,
                width=card_width, height=Inches(0.9).emu,
                text=m['value'],
                font_size=36, font_color=self.cfg.accent, bold=True,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
            )

            # ラベル
            self._add_textbox(
                slide,
                left=x, top=start_y + Inches(1.3).emu,
                width=card_width, height=Inches(0.6).emu,
                text=m['label'],
                font_size=12, font_color=self.cfg.text_secondary,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
            )

        return slide

    def _build_quote(self, slide_data: SlideData):
        """引用スライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        start_y = content_top + Inches(0.5).emu

        # 引用符
        self._add_textbox(
            slide,
            left=Inches(1.5).emu, top=start_y,
            width=Inches(1).emu, height=Inches(0.8).emu,
            text='\u201C',
            font_size=60, font_color=self.cfg.accent, bold=True,
            alignment=PP_ALIGN.LEFT, line_spacing=1.0
        )

        # 引用文
        self._add_textbox(
            slide,
            left=Inches(2.0).emu, top=start_y + Inches(0.3).emu,
            width=Inches(8).emu, height=Inches(2.0).emu,
            text=slide_data.quote_text,
            font_size=18, font_color=self.cfg.text_primary,
            alignment=PP_ALIGN.LEFT, line_spacing=1.6
        )

        # 著者
        if slide_data.quote_author:
            self._add_textbox(
                slide,
                left=Inches(2.0).emu, top=start_y + Inches(2.5).emu,
                width=Inches(8).emu, height=Inches(0.4).emu,
                text=f'\u2014 {slide_data.quote_author}',
                font_size=14, font_color=self.cfg.text_secondary,
                alignment=PP_ALIGN.RIGHT
            )

        return slide

    def _build_faq(self, slide_data: SlideData):
        """FAQスライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        items = slide_data.faq_items
        if not items:
            return slide

        start_y = content_top + Inches(0.2).emu
        item_height = Inches(1.0).emu

        for i, item in enumerate(items):
            y = start_y + i * (item_height + Inches(0.2).emu)

            # Q背景
            self._add_rounded_rect(
                slide,
                left=Inches(0.8).emu, top=y,
                width=Inches(10.4).emu, height=item_height,
                fill_color=self.cfg.bg_secondary
            )

            # Qバッジ
            self._add_rounded_rect(
                slide,
                left=Inches(0.8).emu, top=y,
                width=Inches(0.5).emu, height=Inches(0.4).emu,
                fill_color=self.cfg.accent,
                text='Q', font_size=12, font_color=self.cfg.text_light, bold=True
            )

            # 質問テキスト
            self._add_textbox(
                slide,
                left=Inches(1.5).emu, top=y,
                width=Inches(9.5).emu, height=Inches(0.4).emu,
                text=item['question'],
                font_size=13, font_color=self.cfg.text_primary, bold=True,
                anchor=MSO_ANCHOR.MIDDLE
            )

            # Aバッジ
            self._add_rounded_rect(
                slide,
                left=Inches(0.8).emu, top=y + Inches(0.5).emu,
                width=Inches(0.5).emu, height=Inches(0.4).emu,
                fill_color=self.cfg.secondary,
                text='A', font_size=12, font_color=self.cfg.text_light, bold=True
            )

            # 回答テキスト
            self._add_textbox(
                slide,
                left=Inches(1.5).emu, top=y + Inches(0.5).emu,
                width=Inches(9.5).emu, height=Inches(0.45).emu,
                text=item['answer'],
                font_size=11, font_color=self.cfg.text_secondary,
                anchor=MSO_ANCHOR.MIDDLE
            )

        return slide

    def _build_comparison_table(self, slide_data: SlideData):
        """比較表スライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        table_data = slide_data.table
        if not table_data:
            return slide

        rows = len(table_data.rows) + 1  # ヘッダー含む
        cols = len(table_data.headers)

        start_y = content_top + Inches(0.3).emu
        table_width = Inches(10.8).emu
        table_height = min(Inches(0.45 * rows).emu, SLIDE_HEIGHT - start_y - Inches(0.5).emu)

        tbl = slide.shapes.add_table(
            rows, cols,
            Emu(Inches(0.6).emu), Emu(start_y),
            Emu(table_width), Emu(table_height)
        ).table

        # ヘッダー行
        for j, header in enumerate(table_data.headers):
            cell = tbl.cell(0, j)
            cell.text = header.replace('**', '')
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = self.cfg.font_family
                p.font.bold = True
                p.font.color.rgb = self.cfg.rgb(self.cfg.text_light)
                p.alignment = PP_ALIGN.CENTER
            # ヘッダー背景色
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = self.cfg.rgb(self.cfg.primary)

        # データ行
        for i, row in enumerate(table_data.rows):
            for j, val in enumerate(row):
                if j >= cols:
                    break
                cell = tbl.cell(i + 1, j)
                cell.text = val.replace('**', '')
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.name = self.cfg.font_family
                    p.font.color.rgb = self.cfg.rgb(self.cfg.text_primary)
                    p.alignment = PP_ALIGN.CENTER
                # 交互背景
                cell_fill = cell.fill
                cell_fill.solid()
                if i % 2 == 0:
                    cell_fill.fore_color.rgb = self.cfg.rgb(self.cfg.bg_primary)
                else:
                    cell_fill.fore_color.rgb = self.cfg.rgb(self.cfg.bg_secondary)

        return slide

    def _build_image_with_text(self, slide_data: SlideData):
        """画像+テキスト"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        start_y = content_top + Inches(0.3).emu

        # 画像(左側)
        img_path = slide_data.image_path
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                Emu(Inches(0.6).emu), Emu(start_y),
                Emu(Inches(5).emu), Emu(Inches(3.5).emu)
            )
        else:
            # プレースホルダ表示
            placeholder_path = os.path.join(
                os.path.dirname(__file__), '..', 'assets', 'placeholder.png'
            )
            if os.path.exists(placeholder_path):
                slide.shapes.add_picture(
                    placeholder_path,
                    Emu(Inches(0.6).emu), Emu(start_y),
                    Emu(Inches(5).emu), Emu(Inches(3.5).emu)
                )

        # テキスト(右側)
        self._add_textbox(
            slide,
            left=Inches(6.2).emu, top=start_y,
            width=Inches(5.4).emu, height=Inches(3.5).emu,
            text=slide_data.description,
            font_size=12, font_color=self.cfg.text_primary,
            line_spacing=1.5
        )

        return slide

    def _build_chart(self, slide_data: SlideData):
        """グラフスライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        content_top = self._build_header(slide, slide_data)

        chart_top = content_top + Inches(0.2).emu
        chart_height = SLIDE_HEIGHT - chart_top - Inches(0.5).emu

        self.chart_renderer.render(
            slide, slide_data,
            left=Inches(0.8).emu, top=chart_top,
            width=Inches(10.4).emu, height=chart_height
        )

        return slide

    def _build_cta(self, slide_data: SlideData):
        """行動喚起スライド"""
        slide = self._add_blank_slide()
        self._remove_placeholders(slide)
        self._set_background(slide, self.cfg.bg_secondary)

        # メッセージ(大きく中央)
        self._add_textbox(
            slide,
            left=Inches(1).emu, top=Inches(2.0).emu,
            width=Inches(10).emu, height=Inches(1.0).emu,
            text=slide_data.title,
            font_size=32, font_color=self.cfg.primary, bold=True,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
        )

        # 補足説明
        if slide_data.key_message:
            self._add_textbox(
                slide,
                left=Inches(1).emu, top=Inches(3.2).emu,
                width=Inches(10).emu, height=Inches(0.6).emu,
                text=slide_data.key_message,
                font_size=16, font_color=self.cfg.text_secondary,
                alignment=PP_ALIGN.CENTER
            )

        # ボタン
        if slide_data.button_text:
            self._add_rounded_rect(
                slide,
                left=Inches(4).emu, top=Inches(4.2).emu,
                width=Inches(4).emu, height=Inches(0.7).emu,
                fill_color=self.cfg.accent,
                text=slide_data.button_text,
                font_size=16, font_color=self.cfg.text_light, bold=True
            )

        return slide

    # --- メインビルド ---

    def build_slide(self, slide_data: SlideData):
        """SlideDataに応じたスライドを生成する"""
        layout = slide_data.layout

        if layout == "title":
            return self._build_title(slide_data)
        elif layout == "section":
            return self._build_section(slide_data)
        elif layout == "toc":
            return self._build_toc(slide_data)
        elif layout == "bullet_points":
            return self._build_bullet_points(slide_data)
        elif layout == "numbered_list":
            return self._build_numbered_list(slide_data)
        elif layout == "two_column":
            return self._build_columns(slide_data, 2)
        elif layout == "three_column":
            return self._build_columns(slide_data, 3)
        elif layout == "four_column":
            return self._build_columns(slide_data, 4)
        elif layout == "metrics":
            return self._build_metrics(slide_data)
        elif layout == "quote":
            return self._build_quote(slide_data)
        elif layout == "faq":
            return self._build_faq(slide_data)
        elif layout == "comparison_table":
            return self._build_comparison_table(slide_data)
        elif layout == "image_with_text":
            return self._build_image_with_text(slide_data)
        elif layout == "chart":
            return self._build_chart(slide_data)
        elif layout == "cta":
            return self._build_cta(slide_data)
        else:
            # デフォルト: bullet_points
            return self._build_bullet_points(slide_data)

    def build_all(self, slides: list[SlideData], layout_filter: list[str] | None = None):
        """全スライドを生成する"""
        for slide_data in slides:
            if layout_filter and slide_data.layout not in layout_filter:
                continue
            self.build_slide(slide_data)

    def save(self, output_path: str):
        """PPTXファイルを保存する"""
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        self.prs.save(output_path)
        print(f"Saved: {output_path}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint from slides.md')
    parser.add_argument('--markdown-file', required=True, help='Path to slides.md')
    parser.add_argument('--config', default=None, help='Path to config.json')
    parser.add_argument('--title', default=None, help='Presentation title')
    parser.add_argument('--output-dir', default='.', help='Output directory')
    parser.add_argument('--template', default=None, help='Path to template.pptx')
    parser.add_argument('--layout', default=None, help='Comma-separated layout filter')

    args = parser.parse_args()

    # config読み込み
    config = Config(args.config)

    # テンプレート解決
    template_path = args.template
    if not template_path and config.template_path:
        # config相対パス
        if args.config:
            config_dir = os.path.dirname(os.path.abspath(args.config))
            candidate = os.path.join(config_dir, config.template_path)
            if os.path.exists(candidate):
                template_path = candidate

    # マークダウン解析
    slides = parse_markdown_file(args.markdown_file)
    print(f"Parsed {len(slides)} slides from {args.markdown_file}")

    # レイアウトフィルター
    layout_filter = None
    if args.layout:
        layout_filter = [l.strip() for l in args.layout.split(',')]

    # タイトル決定
    title = args.title
    if not title:
        for s in slides:
            if s.layout == 'title' and s.title:
                title = s.title
                break
    if not title:
        title = 'presentation'

    # ファイル名に使えない文字を除去
    safe_title = re.sub(r'[\\/:*?"<>|]', '', title)

    # ビルド
    builder = PresentationBuilder(config, template_path)
    builder.build_all(slides, layout_filter)

    # 保存
    output_path = os.path.join(args.output_dir, f'{safe_title}.pptx')
    builder.save(output_path)


if __name__ == '__main__':
    main()
